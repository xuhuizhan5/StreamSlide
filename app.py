import os
# import sys
# import importlib
# os.environ["STREAMLIT_SERVER_WATCH_FILE_PATHS"] = "False"
# os.environ["PYTORCH_JIT"] = "0"

# # Apply monkey patch
# if 'torch' in sys.modules:
#     import torch
#     sys.modules['torch'].__path__ = type('', (), {'_path': []})()

# # Create lazy loaders for PyTorch-dependent modules
# class LazyLoader:
#     def __init__(self, module_name):
#         self.module_name = module_name
#         self.module = None
    
#     def __getattr__(self, name):
#         if self.module is None:
#             # Only import when first accessed
#             self.module = importlib.import_module(self.module_name)
            
#             # Apply monkey patch after import
#             if self.module_name == 'docling':
#                 if 'torch' in sys.modules:
#                     import torch
#                     sys.modules['torch'].__path__ = type('', (), {'_path': []})()
        
#         return getattr(self.module, name)

# # Use lazy loading for docling
# docling = LazyLoader('docling')
import io
import json
import time
import uuid
import tempfile
from pathlib import Path
from typing import Dict, List
import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import docling
from docling.document_converter import DocumentConverter, InputFormat, PdfFormatOption
from docling.datamodel.pipeline_options import PdfPipelineOptions    
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
import fitz  # PyMuPDF
import re
import pandas as pd
from pptx.enum.shapes import MSO_SHAPE
import multiprocessing

# Load environment variables
load_dotenv()

# Set page configuration
st.set_page_config(
    page_title="PDF to PowerPoint Generator",
    page_icon="📊",
    layout="wide"
)

# Constants
TEMP_DIR = Path("temp_files")
IMAGE_DIR = TEMP_DIR / "images"
TABLE_DIR = TEMP_DIR / "tables"
OUTPUT_DIR = Path("output")

# Create necessary directories
TEMP_DIR.mkdir(exist_ok=True)
IMAGE_DIR.mkdir(exist_ok=True)
TABLE_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

GEMINI_MODEL = "gemini-1.5-flash"

# Function to setup API key
def setup_gemini_api():
    api_key = st.session_state.get("api_key", os.getenv("GEMINI_API_KEY", ""))
    
    if not api_key:
        st.warning("Gemini API key is not set. Please enter it below.")
        return False
    
    try:
        genai.configure(api_key=api_key)
        # Test connection with model
        model = genai.GenerativeModel(GEMINI_MODEL)
        return True
    except Exception as e:
        st.error(f"Error configuring Gemini API: {str(e)}")
        return False

# def generate_image_caption(image_path: str) -> str:
#     """Generate caption for an image using Gemini Vision model"""
#     try:
#         # Read the image
#         with open(image_path, "rb") as f:
#             image_bytes = f.read()
        
#         # Initialize the Gemini Pro Vision model
#         model = genai.GenerativeModel(GEMINI_MODEL)
        
#         # Prompt for image caption
#         prompt = "Describe this image in 2-3 short sentences. Be concise and focus on what's most important in the image."
        
#         # Generate the caption
#         response = model.generate_content(
#             [prompt, {"mime_type": "image/jpeg", "data": image_bytes}],
#             generation_config={"temperature": 0.2, "max_output_tokens": 256}
#         )
        
#         # Get the caption from the response
#         caption = response.text.strip()
#         return caption
    
#     except Exception as e:
#         print(f"Error generating image caption: {e}")
#         return "Image description not available"

# def generate_table_summary(table_data: Dict) -> str:
#     """Generate a summary of a table using Gemini"""
#     try:
#         # Convert dict to a string representation of a table
#         df = pd.DataFrame.from_dict(table_data)
#         table_str = df.to_markdown(index=False)
        
#         if len(table_str) > 10000:  # If table is too large, truncate
#             table_str = table_str[:10000] + "...(truncated)"
        
#         # Initialize the model
#         model = genai.GenerativeModel(GEMINI_MODEL)
        
#         # Prompt for table summary
#         prompt = """
#         Summarize the following table data. Focus on:
#         1. What this table represents
#         2. Key trends or important information
#         3. Keep your answer under 100 words
        
#         TABLE DATA:
#         """
        
#         # Generate the summary
#         response = model.generate_content(
#             prompt + table_str,
#             generation_config={"temperature": 0.2, "max_output_tokens": 256}
#         )
        
#         # Get the summary from the response
#         summary = response.text.strip()
#         return summary
    
#     except Exception as e:
#         print(f"Error generating table summary: {e}")
#         return "Table summary not available"

def generate_batch_image_captions(image_list: List[dict]) -> Dict[str, str]:
    """Generate captions for multiple images in a single API call"""
    try:
        if not image_list:
            return {}
        
        # Initialize the Gemini Pro Vision model
        model = genai.GenerativeModel(GEMINI_MODEL)
        
        # Maximum batch size to prevent exceeding context limits
        MAX_BATCH_SIZE = 8
        
        # Process in batches
        all_captions = {}
        
        for i in range(0, len(image_list), MAX_BATCH_SIZE):
            batch = image_list[i:i+MAX_BATCH_SIZE]
            
            # Prepare content parts for the model
            content_parts = [
                "Generate concise captions for each of the following images. Each caption should be 1-2 sentences focused on the key elements of the image. Return the results as a valid JSON object with image IDs as keys and captions as values. For any image that appears to be corrupted, nonsensical, or just contains decorative elements like color bars, gradients, or background patterns with no meaningful content, use the value '[NON-INFORMATIVE IMAGE]' instead of attempting to describe it."
            ]
            
            # Add each image to the content parts with its ID
            for idx, img_data in enumerate(batch):
                img_path = img_data["path"]
                img_id = img_data["id"]
                
                with open(img_path, "rb") as f:
                    image_bytes = f.read()
                
                content_parts.append({"mime_type": "image/jpeg", "data": image_bytes})
                content_parts.append(f"Image ID: {img_id}")
            
            # Generate captions for this batch
            response = model.generate_content(
                content_parts,
                generation_config={
                    "temperature": 0.2, 
                    "max_output_tokens": 1024,
                    "response_mime_type": "application/json"
                }
            )
            
            # Parse the response
            try:
                # First try to parse as direct JSON
                batch_captions = json.loads(response.text)
            except json.JSONDecodeError:
                # If direct parsing fails, try to extract JSON from text
                json_match = re.search(r'```json\n(.*?)\n```', response.text, re.DOTALL) or \
                             re.search(r'(\{.*\})', response.text, re.DOTALL)
                
                if json_match:
                    batch_captions = json.loads(json_match.group(1))
                else:
                    print(f"Failed to parse JSON response: {response.text}")
                    batch_captions = {}
            
            # Add this batch's captions to the overall results
            all_captions.update(batch_captions)

        # Delete the items and captions that are not informative
        all_captions = {k: v for k, v in all_captions.items() if v != '[NON-INFORMATIVE IMAGE]'}

        # Save the captions in a json file
        with open(TEMP_DIR / "image_captions.json", "w") as f:
            json.dump(all_captions, f)
        
        return all_captions
    
    except Exception as e:
        print(f"Error in batch image caption generation: {e}")
        return {}

def generate_batch_table_summaries(table_list: List[dict]) -> Dict[str, str]:
    """Generate summaries for multiple tables in a single API call"""
    try:
        if not table_list:
            return {}
        
        # Initialize the model
        model = genai.GenerativeModel(GEMINI_MODEL)
        
        # Maximum batch size
        MAX_BATCH_SIZE = 5
        
        # Process in batches
        all_summaries = {}
        
        for i in range(0, len(table_list), MAX_BATCH_SIZE):
            batch = table_list[i:i+MAX_BATCH_SIZE]
            
            # Build prompt with all table data
            prompt = """
            Generate concise 1-2 sentence descriptions for each of the following tables. Each description should:
            1. Describe what the table represents
            2. Highlight the most important information or pattern
            
            Return the results as a valid JSON object with table IDs as keys and descriptions as values.
            
            TABLES DATA:
            """
            
            for table_data in batch:
                table_id = table_data["id"]
                df = pd.DataFrame.from_dict(table_data["data"])
                table_str = df.to_markdown(index=False)
                
                if len(table_str) > 2000:  # Truncate very large tables
                    table_str = table_str[:2000] + "...(truncated)"
                
                prompt += f"\n\nTable ID: {table_id}\n{table_str}\n"
            
            # Generate summaries for this batch
            response = model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.2, 
                    "max_output_tokens": 1024,
                    "response_mime_type": "application/json"
                }
            )
            
            # Parse the response
            try:
                # First try to parse as direct JSON
                batch_summaries = json.loads(response.text)
            except json.JSONDecodeError:
                # If direct parsing fails, try to extract JSON
                json_match = re.search(r'```json\n(.*?)\n```', response.text, re.DOTALL) or \
                             re.search(r'(\{.*\})', response.text, re.DOTALL)
                
                if json_match:
                    batch_summaries = json.loads(json_match.group(1))
                else:
                    print(f"Failed to parse JSON response: {response.text}")
                    batch_summaries = {}
            
            # Add this batch's summaries to the overall results
            all_summaries.update(batch_summaries)
        
        # Save the summaries in a json file
        with open(TEMP_DIR / "table_summaries.json", "w") as f:
            json.dump(all_summaries, f)

        return all_summaries
    
    except Exception as e:
        print(f"Error in batch table summary generation: {e}")
        return {}

# Process PDF to extract images using PyMuPDF and tables using Docling
def process_pdf(pdf_path: str) -> Dict:
    """Process PDF using PyMuPDF for images and Docling for tables"""
    with st.spinner("Processing PDF..."):
        # Create a progress bar and status message for sub-steps
        progress_placeholder = st.empty()
        status_placeholder = st.empty()
        progress_bar = progress_placeholder.progress(0)
        status_placeholder.text("Initializing PDF processing...")
        
        document_data = {
            "text": "",
            "images": [],
            "tables": []
        }
        
        try:
            # Extract text and images using PyMuPDF
            status_placeholder.text("Step 1/4: Extracting text and images...")
            doc = fitz.open(pdf_path)
            full_text = []
            
            # Track page progress during extraction
            total_pages = len(doc)
            for page_num, page in enumerate(doc):
                # Update progress based on page count
                sub_progress = (page_num + 1) / total_pages * 0.25  # 25% for this step
                progress_bar.progress(sub_progress)
                
                # Extract text
                text = page.get_text()
                full_text.append(text)
                
                # Extract images
                image_list = page.get_images(full=True)
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]  # Image reference
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Get image extension and info
                    image_ext = base_image["ext"]
                    img_id = str(uuid.uuid4())
                    img_path = str(IMAGE_DIR / f"image_{img_id}.{image_ext}")
                    
                    # Save image
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # Generate a placeholder for image caption (will be filled later)
                    images_data = {
                        "id": img_id,
                        "path": img_path,
                        "page": page_num + 1,
                        "caption": "",
                        "description": ""
                    }
                    
                    document_data["images"].append(images_data)
            
            # Combine all text
            document_data["text"] = "\n".join(full_text)
            
            # Update progress and status for table extraction
            progress_bar.progress(0.25)  # 25% completion
            status_placeholder.text("Step 2/4: Extracting tables...")
            
            # Now extract tables using Docling
            try:
                # Configure pipeline options for tables only
                pipeline_options = PdfPipelineOptions()
                pipeline_options.do_picture_description = False  # Don't extract images with Docling
                pipeline_options.generate_page_images = False
                pipeline_options.generate_picture_images = False
                
                # Initialize the DocumentConverter
                converter = DocumentConverter(format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                })
                
                # Convert the document
                result = converter.convert(pdf_path)
                document = result.document
                
                # Process tables
                table_count = sum(1 for item, _ in document.iterate_items() if hasattr(item, 'export_to_dataframe'))
                processed_tables = 0
                
                for item, _ in document.iterate_items():
                    if hasattr(item, 'export_to_dataframe'):
                        table_id = str(uuid.uuid4())
                        table_df = item.export_to_dataframe()
                        table_path = str(TABLE_DIR / f"table_{table_id}.csv")

                        table_df.to_csv(table_path, index=False)
                        
                        table_data = {
                            "id": table_id,
                            "path": table_path,
                            "caption": item.caption if hasattr(item, 'caption') else "",
                            "data": table_df.to_dict()
                        }
                        
                        document_data["tables"].append(table_data)
                        
                        # Update progress for table extraction
                        processed_tables += 1
                        if table_count > 0:
                            sub_progress = 0.25 + (processed_tables / table_count * 0.25)  # 25% to 50%
                            progress_bar.progress(sub_progress)
                
            except Exception as e:
                st.warning(f"Error extracting tables with Docling: {str(e)}")
                st.info("Continuing without table extraction...")
                progress_bar.progress(0.5)  # Move to 50% if table extraction fails
                
            # Update progress for image caption generation
            progress_bar.progress(0.5)  # 50% completion
            
            # Generate captions for images using Gemini in batches
            if document_data["images"]:
                status_placeholder.text(f"Step 3/4: Generating captions for {len(document_data['images'])} images...")
                
                # Use batch processing for images
                image_captions = generate_batch_image_captions(document_data["images"])
                
                # Assign captions to images
                for img_data in document_data["images"]:
                    img_id = img_data["id"]
                    if img_id in image_captions:
                        img_data["caption"] = image_captions[img_id]
                        img_data["description"] = image_captions[img_id]
                    else:
                        img_data["caption"] = "Image description not available"
                        img_data["description"] = "Image description not available"
                
                progress_bar.progress(0.75)  # 75% completion
            else:
                # Skip this step if no images
                progress_bar.progress(0.75)
            
            # Generate summaries for tables in batches
            status_placeholder.text("Step 4/4: Generating summaries for tables...")
            
            if document_data["tables"]:
                # Filter tables that need summaries
                tables_needing_summaries = [
                    table for table in document_data["tables"]
                    if not table.get("caption") or table.get("caption") == ""
                ]
                
                if tables_needing_summaries:
                    # Use batch processing for tables
                    table_summaries = generate_batch_table_summaries(tables_needing_summaries)
                    
                    # Assign summaries to tables
                    for table_data in document_data["tables"]:
                        table_id = table_data["id"]
                        if table_id in table_summaries:
                            table_data["caption"] = table_summaries[table_id]
            
            # Complete progress
            progress_bar.progress(1.0)
            status_placeholder.text("PDF processing complete!")
            
            return document_data
            
        except Exception as e:
            st.error(f"Error processing PDF: {str(e)}")
            return document_data

# Function to optimize prompt based on presentation time
def optimize_prompt_for_presentation_time(basic_prompt, document_content, presentation_time):
    """Use multi-step prompt optimization based on presentation time"""
    with st.spinner("Optimizing presentation structure..."):
        try:
            model = genai.GenerativeModel(GEMINI_MODEL)
            
            # First step: Analyze content and time requirements
            system_prompt = """
            You are an expert presentation coach specializing in time management for presentations.
            Based on the presentation topic and available time, determine:
            1. Optimal number of slides
            2. Content density per slide
            3. Time allocation strategy
            
            Provide specific, actionable guidance on how to structure this presentation.
            """
            
            # Extract a summary of the document content
            if isinstance(document_content, dict):
                doc_summary = document_content.get("text", "")
                if len(doc_summary) > 5000:
                    doc_summary = doc_summary[:5000] + "..."
            else:
                doc_summary = str(document_content)[:5000] + "..."
            
            time_analysis_prompt = f"""
            Basic presentation request: {basic_prompt}
            
            Document summary: {doc_summary}
            
            Available presentation time: {presentation_time} minutes
            
            Please analyze and provide structured guidance for this presentation.
            """
            
            response = model.generate_content(
                [system_prompt, time_analysis_prompt],
                generation_config={"temperature": 0.2}
            )
            
            time_analysis = response.text
            
            # Second step: Create optimized prompt with presentation structure
            system_prompt_2 = """
            You are an expert presentation designer who optimizes content for specific time constraints.
            Convert the provided analysis into a detailed, structured prompt that will guide the creation 
            of a well-paced presentation matching the time requirements.
            
            Your output should be a comprehensive prompt that specifies:
            1. Exact number of slides to create
            2. Specific slide types and their order
            3. Content density guidelines for each section
            4. Transitions and timing suggestions
            """
            
            optimization_prompt = f"""
            Original request: {basic_prompt}
            
            Time analysis: {time_analysis}
            
            Available presentation time: {presentation_time} minutes
            
            Document summary: {doc_summary}
            
            Create an optimized, structured prompt for generating this presentation.
            """
            
            response = model.generate_content(
                [system_prompt_2, optimization_prompt],
                generation_config={"temperature": 0.2}
            )
            
            optimized_prompt = response.text
            
            # Combine with original prompt for final result
            final_prompt = f"""
            {basic_prompt}
            
            PRESENTATION TIME: {presentation_time} minutes
            
            STRUCTURAL GUIDANCE:
            {optimized_prompt}
            """
            
            return final_prompt
            
        except Exception as e:
            st.warning(f"Error optimizing prompt: {str(e)}")
            return basic_prompt  # Fallback to original prompt if optimization fails

# Function to generate presentation content using Gemini
def generate_presentation_content(document_data: Dict, user_prompt: str) -> Dict:
    """Generate presentation content using Gemini with document data"""
    with st.spinner("Generating presentation content..."):
        try:
            # Extract document text
            document_text = document_data.get("text", "")
            
            # Truncate if too long
            if len(document_text) > 100000:
                document_text = document_text[:100000] + "..."
            
            # Prepare image references - add additional info about each image
            image_references = []
            for img in document_data.get("images", []):
                image_references.append({
                    "id": img["id"],
                    "path": img["path"],
                    "caption": img["caption"],
                    "page": img.get("page", "unknown")
                })
            
            # Prepare table references with summaries
            table_references = []
            for table in document_data.get("tables", []):
                table_references.append({
                    "id": table["id"],
                    "path": table["path"],
                    "caption": table["caption"],
                    "summary": table.get("caption", "Table data")
                })
            
            # Prepare prompt for Gemini
            system_prompt = """
            You are an expert presentation designer specialized in creating professional research presentations.

            Create a research-focused presentation with appropriate slide types based on the paper's content:
            1. Title Slide: Paper title, authors, affiliations
            2. Introduction/Agenda: Research question, objectives, and brief overview
            3. Background/Literature: Key related work and context
            4. Methods/Approach: Research methodology
            5. Results: Key findings with supporting data
            6. Discussion: Interpretation of results
            7. Conclusion: Summary of contributions and implications

            IMPORTANT CONTENT GUIDELINES FOR RESEARCH PRESENTATIONS:
            - Focus on logical flow of research narrative
            - Include tables WITHIN content slides when they directly support the point being made
            - Use dedicated table slides only for complex tables that need detailed explanation
            - Create image slides when visuals help explain concepts or results
            - Balance text and visual elements based on the research content
            - Use section slides to clearly delineate major parts of the research

            For slide type distribution in research presentations:
            - 1 title slide
            - 1 introduction/agenda slide
            - 2-3 background/related work slides
            - 2-4 methodology slides
            - 3-6 results slides (with tables and figures as appropriate)
            - 2-3 discussion slides
            - 1 conclusion slide

            PRESENTATION FORMATTING GUIDELINES:
            - Keep slides focused on key points - avoid overcrowding
            - Use hierarchical structure for complex information
            - Include tables when they directly support the narrative
            - Select relevant images that illustrate important concepts
            - Ensure all slides have clear, descriptive titles
            - Balance text with visual elements

            IMPORTANT SLIDE CONTENT RULES:
            - Tables should be included where they logically fit in the research narrative
            - NO MORE THAN 1 TABLE per slide to ensure readability
            - NO MORE THAN 4 IMAGES per slide
            - Never mix tables and images on the same slide
            - Use appropriate slide types based on content purpose, not arbitrary rules

            IMPORTANT: Your response must be valid JSON with no syntax errors. Double-check all commas, 
            brackets, and ensure all strings are properly quoted. Do not include any explanatory text 
            outside the JSON structure.

            Generate a JSON structure with:
            {
                "title": "Research Paper Title",
                "subtitle": "Authors and Affiliations",
                "theme": {
                    "primary_color": "#3a86ff",  // Main brand color
                    "secondary_color": "#8338ec", // Secondary accent color
                    "text_color": "#333333",     // Primary text color
                    "background_color": "#ffffff", // Slide background
                    "accent_color": "#ff006e"    // Highlight color
                },
                "slides": [
                    {
                        "slide_type": "title_slide|agenda_slide|section_slide|content_slide|image_slide|table_slide|conclusion_slide",
                        "title": "Slide Title",
                        "subtitle": "Optional Slide Subtitle",
                        "content": ["Point 1", "Point 2", "Point 3"],
                        "notes": "Speaker notes for this slide",
                        "images": ["image_path_1", "image_path_2"],
                        "tables": ["table_path_1"]
                    }
                ]
            }

            Include a theme object with professional color choices appropriate for academic/research presentations.
            """
            
            # Include image and table information in the prompt
            user_message = f"""
            {user_prompt}
            
            DOCUMENT CONTENT:
            {document_text}
            
            Available images (reference these by path when creating slides):
            {json.dumps(image_references, indent=2)}
            
            Available tables (reference these by path when creating slides):
            {json.dumps(table_references, indent=2)}
            
            Create a cohesive presentation with logical flow, ensuring:
            1. Clear transitions between topics
            2. Consistent structure on each slide (title + content)
            3. Proper distribution of content (not too dense, not too sparse)
            4. Strategic use of available images and tables where relevant
            5. Professional, engaging language
            
            When including images or tables in slides, select them based on their captions and relevance to the slide content.
            
            Include a title slide, agenda slide, and conclusion slide in your structure.
            
            IMPORTANT: Your response must be valid, properly formatted JSON with no syntax errors.
            Ensure all arrays and objects have proper delimiters (commas between items).
            """
            
            model = genai.GenerativeModel(GEMINI_MODEL)

            response = model.generate_content(
                [system_prompt, user_message],
                generation_config={"temperature": 0.2, "max_output_tokens": 8192}
            )
            
            # Extract JSON from response
            response_text = response.text
            
            # First try to extract JSON with code block markers
            json_match = re.search(r'```json\n(.*?)\n```', response_text, re.DOTALL)
            
            if json_match:
                json_str = json_match.group(1)
            else:
                # Try to extract JSON without code block markers
                json_match = re.search(r'(\{.*\})', response_text, re.DOTALL)
                if json_match:
                    json_str = json_match.group(1)
                else:
                    raise ValueError("Could not extract JSON from response")
            
            # Try to parse the JSON with error handling
            try:
                presentation_content = json.loads(json_str)
            except json.JSONDecodeError as e:
                st.warning(f"Initial JSON parsing failed: {str(e)}")
                
                # Log the problematic JSON for debugging
                with open(TEMP_DIR / "json_error.txt", "w") as f:
                    f.write(f"Error: {str(e)}\n\n")
                    f.write(json_str)
                
                # Try to fix common JSON syntax errors and retry
                fixed_json = fix_json_syntax(json_str)
                try:
                    presentation_content = json.loads(fixed_json)
                    st.info("Fixed JSON syntax errors automatically")
                except json.JSONDecodeError:
                    # If still failing, request a cleaner JSON response
                    st.warning("Retrying with explicit request for clean JSON...")
                    
                    repair_prompt = """
                    The JSON you provided has syntax errors. Please generate a corrected version
                    with proper JSON syntax. Common issues include:
                    1. Missing commas between array or object items
                    2. Trailing commas at the end of arrays or objects
                    3. Unquoted property names
                    4. Mismatched quotes or brackets
                    
                    Respond ONLY with the corrected JSON structure, nothing else.
                    """
                    
                    repair_response = model.generate_content([repair_prompt, json_str])
                    repair_text = repair_response.text
                    
                    # Try to extract and parse the repaired JSON
                    json_match = re.search(r'```json\n(.*?)\n```', repair_text, re.DOTALL)
                    if json_match:
                        repaired_json = json_match.group(1)
                    else:
                        repaired_json = re.search(r'(\{.*\})', repair_text, re.DOTALL).group(1)
                    
                    presentation_content = json.loads(repaired_json)
            
            # Save presentation content
            with open(TEMP_DIR / "presentation_content.json", "w") as f:
                json.dump(presentation_content, f, indent=2)
            
            return presentation_content
            
        except Exception as e:
            st.error(f"Error generating presentation content: {str(e)}")
            return {}

def fix_json_syntax(json_str):
    """Attempt to fix common JSON syntax errors"""
    # Fix missing commas between objects in arrays
    fixed = re.sub(r'}\s*{', '},{', json_str)
    
    # Fix missing commas between array elements
    fixed = re.sub(r'"\s*"', '","', fixed)
    
    # Fix trailing commas in arrays and objects
    fixed = re.sub(r',\s*}', '}', fixed)
    fixed = re.sub(r',\s*\]', ']', fixed)
    
    # Fix unquoted property names
    fixed = re.sub(r'(\w+):', r'"\1":', fixed)
    
    # Remove any extra text before/after the JSON object
    json_obj_match = re.search(r'({.*})', fixed, re.DOTALL)
    if json_obj_match:
        fixed = json_obj_match.group(1)
    
    return fixed

# Function to convert hex color to RGB
def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Create PowerPoint presentation leveraging python-pptx more directly
def create_powerpoint(presentation_content: Dict, document_data: Dict) -> str:
    """Create PowerPoint file from content with professional styling"""
    with st.spinner("Creating PowerPoint presentation..."):
        try:
            # Create presentation
            prs = Presentation()
            
            # Get theme colors from presentation content
            theme = presentation_content.get("theme", {})
            primary_color = theme.get("primary_color", "#3a86ff")
            secondary_color = theme.get("secondary_color", "#8338ec")
            text_color = theme.get("text_color", "#333333")
            bg_color = theme.get("background_color", "#ffffff")
            
            # Convert hex colors to RGB for use in the presentation
            primary_rgb = hex_to_rgb(primary_color)
            secondary_rgb = hex_to_rgb(secondary_color)
            text_rgb = hex_to_rgb(text_color)
            bg_rgb = hex_to_rgb(bg_color)
            
            # Create lookup dictionaries for images and tables by path
            image_lookup = {img["path"]: img for img in document_data.get("images", [])}
            table_lookup = {table["path"]: table for table in document_data.get("tables", [])}
            
            # Process slides based on content
            for slide_content in presentation_content.get("slides", []):
                slide_type = slide_content.get("slide_type", "content_slide")
                
                # Choose appropriate slide layout based on slide type
                if slide_type == "title_slide":
                    slide_layout = prs.slide_layouts[0]  # Title slide layout
                    
                elif slide_type == "agenda_slide":
                    # Use a content layout for agenda slides
                    slide_layout = prs.slide_layouts[1]  # Title and content layout
                elif slide_type == "section_slide":
                    # Find a title-only layout if available
                    slide_layout = prs.slide_layouts[5]
                elif slide_type == "image_slide":
                    slide_layout = prs.slide_layouts[5]  # Title and content layout
                elif slide_type == "table_slide":
                    slide_layout = prs.slide_layouts[5]  # Title and content layout
                else:  # content_slide or conclusion_slide
                    slide_layout = prs.slide_layouts[1]  # Title and content layout
                
                # Add slide with the chosen layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_rgb)
                
                # Add title
                if hasattr(slide.shapes, 'title') and slide.shapes.title:
                    title_shape = slide.shapes.title
                    title_shape.text = slide_content.get("title", "")
                    
                    # Style title
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(*primary_rgb if slide_type != "section_slide" else (255, 255, 255))
                            run.font.bold = True
                            run.font.size = Pt(40 if slide_type == "section_slide" else 32)
                
                # Add subtitle for title slide
                if slide_type == "title_slide" and "subtitle" in slide_content:
                    for shape in slide.placeholders:
                        if shape.placeholder_format.type == 2:  # Subtitle placeholder
                            shape.text = slide_content.get("subtitle", "")
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(*text_rgb)
                                    run.font.size = Pt(28)
                                    run.font.italic = True
                
                # Special handling for agenda slides
                if slide_type == "agenda_slide":
                    content = slide_content.get("content", [])
                    if content:
                        # Find a content placeholder or create a textbox
                        shape = None
                        for placeholder in slide.placeholders:
                            if placeholder.placeholder_format.type not in (1, 2):  # Not title or subtitle
                                shape = placeholder
                                break
                        
                        if not shape:
                            # Create textbox if no placeholder is available
                            left = Inches(1)
                            top = Inches(2)
                            width = Inches(8)
                            height = Inches(5)
                            shape = slide.shapes.add_textbox(left, top, width, height)
                        
                        # Clear any existing text
                        text_frame = shape.text_frame
                        text_frame.clear()
                        
                        # Process each agenda item
                        for i, item in enumerate(content):
                            # Determine if it's a main section or subsection
                            is_subsection = item.startswith("  ") or item.startswith("- ")
                            level = 1 if is_subsection else 0
                            
                            # Clean up the item text
                            item_text = item.strip()
                            if item_text.startswith("- ") or item_text.startswith("• "):
                                item_text = item_text[2:].strip()
                            
                            # Add paragraph with proper indentation
                            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                            p.text = item_text
                            p.level = level
                            
                            # Style based on level
                            for run in p.runs:
                                run.font.size = Pt(24 if level == 0 else 20)
                                run.font.color.rgb = RGBColor(*primary_rgb if level == 0 else text_rgb)
                                run.font.bold = (level == 0)
                            
                            # Add spacing between main sections
                            if level == 0 and i > 0:
                                p.space_before = Pt(18)
                
                # Handle content slides (not agenda slides)
                elif slide_type != "title_slide":
                    # Add a subtitle
                    if "subtitle" in slide_content:
                        for shape in slide.placeholders:
                            if shape.placeholder_format.type == 2:  # Subtitle placeholder
                                shape.text = slide_content.get("subtitle", "")
                                break
                    
                    # Add content (bullet points)
                    content = slide_content.get("content", [])
                    if content:
                        # Find a content placeholder
                        content_placeholder = None
                        for shape in slide.placeholders:
                            if shape.placeholder_format.type not in (1, 2):  # Not title or subtitle
                                content_placeholder = shape
                                break
                        
                        if content_placeholder:
                            # Use the placeholder for bullet points
                            tf = content_placeholder.text_frame
                            tf.clear()
                            
                            for i, point in enumerate(content):
                                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                                p.text = point
                                p.level = 0
                                
                                for run in p.runs:
                                    run.font.size = Pt(18 if slide_type != "conclusion_slide" else 20)
                                    run.font.color.rgb = RGBColor(*text_rgb)
                        else:
                            # Create textbox if no placeholder is available
                            left = Inches(1)
                            top = Inches(2)
                            width = Inches(8)
                            height = Inches(5)
                            
                            txtbox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txtbox.text_frame
                            
                            for i, point in enumerate(content):
                                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                                p.text = point
                                
                                for run in p.runs:
                                    run.font.size = Pt(18 if slide_type != "conclusion_slide" else 20)
                                    run.font.color.rgb = RGBColor(*text_rgb)
                
                # Add images with captions
                images = slide_content.get("images", None)
                if images:
                    # Position logic based on number of images and slide type
                    img_positions = []
                    
                    # Use different positioning for dedicated image slides
                    if slide_type == "image_slide":
                        # For image slides, make images larger and more centered
                        if len(images) == 1:
                            # Single image gets more screen real estate
                            img_positions = [(Inches(1.5), Inches(2.2), Inches(7), Inches(0.6))]
                        elif len(images) == 2:
                            img_positions = [
                                (Inches(1), Inches(2.2), Inches(4), Inches(0.6)),
                                (Inches(5), Inches(2.2), Inches(4), Inches(0.6))
                            ]
                        elif len(images) <= 4:
                            img_positions = [
                                (Inches(1), Inches(2.2), Inches(4), Inches(0.6)),
                                (Inches(5), Inches(2.2), Inches(4), Inches(0.6)),
                                (Inches(1), Inches(4.5), Inches(4), Inches(0.6)),
                                (Inches(5), Inches(4.5), Inches(4), Inches(0.6))
                            ]
                    else:
                        # Standard image positioning for content slides with images
                        if len(images) == 1:
                            img_positions = [(Inches(2), Inches(3.5), Inches(6), Inches(0.6))]
                        elif len(images) == 2:
                            img_positions = [
                                (Inches(1), Inches(3.5), Inches(3.5), Inches(0.6)),
                                (Inches(5), Inches(3.5), Inches(3.5), Inches(0.6))
                            ]
                        elif len(images) <= 4:
                            img_positions = [
                                (Inches(1), Inches(3.5), Inches(3.5), Inches(0.6)),
                                (Inches(5), Inches(3.5), Inches(3.5), Inches(0.6)),
                                (Inches(1), Inches(5.5), Inches(3.5), Inches(0.6)),
                                (Inches(5), Inches(5.5), Inches(3.5), Inches(0.6))
                            ]
                    
                    # Limit to maximum 4 images per slide
                    images = images[:min(len(images), 4)]
                    
                    # Add images to slide
                    successfully_added_images = []
                    for i, img_path in enumerate(images[:min(len(images), len(img_positions))]):
                        if os.path.exists(img_path):
                            left, top, width, _ = img_positions[i]
                            
                            # Add the image
                            slide.shapes.add_picture(img_path, left, top, width=width)
                            successfully_added_images.append(img_path)
                    
                    # Add a single consolidated caption for all images on the slide - ONLY if images were added
                    if successfully_added_images:
                        # Create a caption that combines information or uses a general description
                        combined_caption = "Images: "
                        image_captions = []
                        
                        for img_path in successfully_added_images:
                            if img_path in image_lookup and image_lookup[img_path].get("caption"):
                                image_captions.append(image_lookup[img_path]["caption"])
                        
                        # If we have captions, create a consolidated one
                        if image_captions:
                            if len(image_captions) == 1:
                                combined_caption = image_captions[0]
                            else:
                                # For multiple images, use a more intelligent consolidated caption
                                if len(image_captions) > 1:
                                    try:
                                        # Initialize the model
                                        model = genai.GenerativeModel(GEMINI_MODEL)
                                        
                                        # Create a prompt for generating a consolidated caption
                                        prompt = f"""
                                        Create a brief, cohesive caption (under 25 words) that describes this collection of images together.
                                        Focus on their relationship, common themes, or progression.
                                        
                                        Individual image captions:
                                        {chr(10).join([f"- {caption}" for caption in image_captions])}
                                        """
                                        
                                        response = model.generate_content(
                                            prompt,
                                            generation_config={"temperature": 0.2, "max_output_tokens": 100}
                                        )
                                        
                                        combined_caption = response.text.strip()
                                        
                                        # Add fallback in case of empty response
                                        if not combined_caption:
                                            combined_caption = "Collection of related images showing " + ", ".join(image_captions[:2]) + "..."
                                    except Exception as e:
                                        print(f"Error generating combined caption: {e}")
                                        combined_caption = "Collection of related images for this slide"
                        
                        # Add the combined caption below the images
                        caption_box = slide.shapes.add_textbox(
                            Inches(1), Inches(6.5), 
                            Inches(8), Inches(1.0)  # Increased height to accommodate multiple lines
                        )
                        caption_frame = caption_box.text_frame
                        caption_frame.text = combined_caption
                        
                        # Enable word wrapping for multi-line captions
                        caption_frame.word_wrap = True
                        
                        # Style caption
                        for paragraph in caption_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(10)
                                run.font.italic = True
                                run.font.color.rgb = RGBColor(*text_rgb)
                
                # Add tables with captions if present in the slide (limit to 1 table per slide)
                tables = slide_content.get("tables", [])
                if tables:
                    # Take only the first table
                    table_path = tables[0]
                    
                    if table_path in table_lookup:
                        table_data = table_lookup[table_path]
                        
                        # Get table DataFrame
                        try:
                            df = pd.read_csv(table_path)
                            
                            # Create simplified table reference in slides
                            rows, cols = df.shape
                            max_rows = min(rows, 10)  # Limit to 10 rows max for presentation
                            max_cols = min(cols, 6)   # Limit to 6 columns
                            
                            # Position table
                            left = Inches(1)
                            top = Inches(2.5)
                            width = Inches(8)
                            height = Inches(3)
                            
                            # Add table to slide
                            shape_table = slide.shapes.add_table(
                                max_rows + 1,  # +1 for header
                                max_cols,
                                left, top, width, height
                            )
                            
                            tbl = shape_table.table
                            
                            # Add headers
                            for j, col_name in enumerate(df.columns[:max_cols]):
                                cell = tbl.cell(0, j)
                                cell.text = str(col_name)
                                # Style headers
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.color.rgb = RGBColor(*primary_rgb)
                            
                            # Add data
                            for row_idx in range(max_rows):
                                for col_idx in range(max_cols):
                                    if row_idx < rows and col_idx < cols:
                                        cell = tbl.cell(row_idx + 1, col_idx)  # +1 for header row
                                        cell.text = str(df.iloc[row_idx, col_idx])
                            
                            # Add table caption if available
                            if table_data.get("caption"):
                                caption = table_data["caption"]
                                caption_box = slide.shapes.add_textbox(
                                    left, top + Inches(3.2), 
                                    width, Inches(1.0)  # Increased height to accommodate multiple lines
                                )
                                caption_frame = caption_box.text_frame
                                caption_frame.text = f"Table: {caption}"
                                
                                # Enable word wrapping for multi-line captions
                                caption_frame.word_wrap = True
                                
                                # Style caption
                                for paragraph in caption_frame.paragraphs:
                                    paragraph.alignment = PP_ALIGN.CENTER
                                    for run in paragraph.runs:
                                        run.font.size = Pt(10)
                                        run.font.italic = True
                                        run.font.color.rgb = RGBColor(*text_rgb)
                        
                        except Exception as table_err:
                            print(f"Error adding table to slide: {table_err}")
                
                # Add speaker notes if available
                if "notes" in slide_content and slide_content["notes"]:
                    slide.notes_slide.notes_text_frame.text = slide_content["notes"]
            
            # Save the presentation with timestamp
            timestamp = time.strftime("%Y%m%d_%H%M%S")
            output_filename = f"presentation_{timestamp}.pptx"
            output_path = str(OUTPUT_DIR / output_filename)
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            st.error(f"Error creating PowerPoint: {str(e)}")
            return ""

# def create_powerpoint(presentation_content: Dict, document_data: Dict) -> str:
#     """Create PowerPoint file from content with professional styling"""
#     with st.spinner("Creating PowerPoint presentation..."):
#         try:
#             # Create presentation
#             prs = Presentation()
            
#             # Get theme colors from presentation content
#             theme = presentation_content.get("theme", {})
#             primary_color = theme.get("primary_color", "#3a86ff")
#             secondary_color = theme.get("secondary_color", "#8338ec")
#             text_color = theme.get("text_color", "#333333")
#             bg_color = theme.get("background_color", "#ffffff")
            
#             # Convert hex colors to RGB for use in the presentation
#             primary_rgb = hex_to_rgb(primary_color)
#             secondary_rgb = hex_to_rgb(secondary_color)
#             text_rgb = hex_to_rgb(text_color)
#             bg_rgb = hex_to_rgb(bg_color)
            
#             # Create lookup dictionaries for images and tables by path
#             image_lookup = {img["path"]: img for img in document_data.get("images", [])}
#             table_lookup = {table["path"]: table for table in document_data.get("tables", [])}
            
#             # Process slides based on content
#             for slide_content in presentation_content.get("slides", []):
#                 slide_type = slide_content.get("slide_type", "content_slide")
                
#                 # Choose appropriate slide layout based on slide type
#                 if slide_type == "title_slide":
#                     slide_layout = prs.slide_layouts[0]  # Title slide layout
                    
#                 elif slide_type == "agenda_slide":
#                     slide_layout = prs.slide_layouts[1]  # Title and content layout
#                 elif slide_type == "section_slide":
#                     slide_layout = prs.slide_layouts[5]
#                 elif slide_type == "image_slide":
#                     slide_layout = prs.slide_layouts[5]  # Title and content layout
#                 elif slide_type == "table_slide":
#                     slide_layout = prs.slide_layouts[5]  # Title and content layout
#                 else:  # content_slide or conclusion_slide
#                     slide_layout = prs.slide_layouts[1]  # Title and content layout
                
#                 # Add slide with the chosen layout
#                 slide = prs.slides.add_slide(slide_layout)
                
#                 # Set background color
#                 background = slide.background
#                 fill = background.fill
#                 fill.solid()
#                 fill.fore_color.rgb = RGBColor(*bg_rgb)
                
#                 # Add title
#                 if hasattr(slide.shapes, 'title') and slide.shapes.title:
#                     title_shape = slide.shapes.title
#                     title_shape.text = slide_content.get("title", "")
                    
#                     # Style title
#                     for paragraph in title_shape.text_frame.paragraphs:
#                         for run in paragraph.runs:
#                             run.font.color.rgb = RGBColor(*primary_rgb if slide_type != "section_slide" else (255, 255, 255))
#                             run.font.bold = True
#                             run.font.size = Pt(40 if slide_type == "section_slide" else 32)
                
#                 # Add subtitle for title slide
#                 if slide_type == "title_slide" and "subtitle" in slide_content:
#                     for shape in slide.placeholders:
#                         if shape.placeholder_format.type == 2:  # Subtitle placeholder
#                             shape.text = slide_content.get("subtitle", "")
#                             for paragraph in shape.text_frame.paragraphs:
#                                 for run in paragraph.runs:
#                                     run.font.color.rgb = RGBColor(*text_rgb)
#                                     run.font.size = Pt(28)
#                                     run.font.italic = True
                
#                 # Special handling for agenda slides
#                 if slide_type == "agenda_slide":
#                     content = slide_content.get("content", [])
#                     if content:
#                         # Find a content placeholder or create a textbox
#                         shape = None
#                         for placeholder in slide.placeholders:
#                             if placeholder.placeholder_format.type not in (1, 2):  # Not title or subtitle
#                                 shape = placeholder
#                                 break
                        
#                         if not shape:
#                             # Create textbox if no placeholder is available
#                             left = Inches(1)
#                             top = Inches(2)
#                             width = Inches(8)
#                             height = Inches(5)
#                             shape = slide.shapes.add_textbox(left, top, width, height)
                        
#                         # Clear any existing text
#                         text_frame = shape.text_frame
#                         text_frame.clear()
                        
#                         # Process each agenda item
#                         for i, item in enumerate(content):
#                             # Determine if it's a main section or subsection
#                             is_subsection = item.startswith("  ") or item.startswith("- ")
#                             level = 1 if is_subsection else 0
                            
#                             # Clean up the item text
#                             item_text = item.strip()
#                             if item_text.startswith("- ") or item_text.startswith("• "):
#                                 item_text = item_text[2:].strip()
                            
#                             # Add paragraph with proper indentation
#                             p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
#                             p.text = item_text
#                             p.level = level
                            
#                             # Style based on level
#                             for run in p.runs:
#                                 run.font.size = Pt(24 if level == 0 else 20)
#                                 run.font.color.rgb = RGBColor(*primary_rgb if level == 0 else text_rgb)
#                                 run.font.bold = (level == 0)
                            
#                             # Add spacing between main sections
#                             if level == 0 and i > 0:
#                                 p.space_before = Pt(18)
                
#                 # Handle content slides (not agenda slides)
#                 elif slide_type != "title_slide":
#                     # Add a subtitle
#                     if "subtitle" in slide_content:
#                         for shape in slide.placeholders:
#                             if shape.placeholder_format.type == 2:  # Subtitle placeholder
#                                 shape.text = slide_content.get("subtitle", "")
#                                 break
                    
#                     # Add content (bullet points)
#                     content = slide_content.get("content", [])
#                     if content:
#                         # Find a content placeholder
#                         content_placeholder = None
#                         for shape in slide.placeholders:
#                             if shape.placeholder_format.type not in (1, 2):  # Not title or subtitle
#                                 content_placeholder = shape
#                                 break
                        
#                         if content_placeholder:
#                             # Use the placeholder for bullet points
#                             tf = content_placeholder.text_frame
#                             tf.clear()
                            
#                             for i, point in enumerate(content):
#                                 p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
#                                 p.text = point
#                                 p.level = 0
                                
#                                 for run in p.runs:
#                                     run.font.size = Pt(18 if slide_type != "conclusion_slide" else 20)
#                                     run.font.color.rgb = RGBColor(*text_rgb)
#                         else:
#                             # Create textbox if no placeholder is available
#                             left = Inches(1)
#                             top = Inches(2)
#                             width = Inches(8)
#                             height = Inches(5)
                            
#                             txtbox = slide.shapes.add_textbox(left, top, width, height)
#                             tf = txtbox.text_frame
                            
#                             for i, point in enumerate(content):
#                                 p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
#                                 p.text = point
                                
#                                 for run in p.runs:
#                                     run.font.size = Pt(18 if slide_type != "conclusion_slide" else 20)
#                                     run.font.color.rgb = RGBColor(*text_rgb)
                
#                 # Add images with captions
#                 images = slide_content.get("images", None)
#                 if images and slide_type != "title_slide":  # Ensure no images on title slides
#                     # Position logic based on number of images and slide type
#                     img_positions = []
                    
#                     # Use different positioning for dedicated image slides
#                     if slide_type == "image_slide":
#                         # For image slides, make images larger and more centered
#                         if len(images) == 1:
#                             # Single image gets more screen real estate
#                             img_positions = [(Inches(1.5), Inches(2.2), Inches(7), Inches(0.6))]
#                         elif len(images) == 2:
#                             img_positions = [
#                                 (Inches(1), Inches(2.2), Inches(4), Inches(0.6)),
#                                 (Inches(5), Inches(2.2), Inches(4), Inches(0.6))
#                             ]
#                         elif len(images) <= 4:
#                             img_positions = [
#                                 (Inches(1), Inches(2.2), Inches(4), Inches(0.6)),
#                                 (Inches(5), Inches(2.2), Inches(4), Inches(0.6)),
#                                 (Inches(1), Inches(4.5), Inches(4), Inches(0.6)),
#                                 (Inches(5), Inches(4.5), Inches(4), Inches(0.6))
#                             ]
#                     else:
#                         # Standard image positioning for content slides with images
#                         if len(images) == 1:
#                             img_positions = [(Inches(2), Inches(3.5), Inches(6), Inches(0.6))]
#                         elif len(images) == 2:
#                             img_positions = [
#                                 (Inches(1), Inches(3.5), Inches(3.5), Inches(0.6)),
#                                 (Inches(5), Inches(3.5), Inches(3.5), Inches(0.6))
#                             ]
#                         elif len(images) <= 4:
#                             img_positions = [
#                                 (Inches(1), Inches(3.5), Inches(3.5), Inches(0.6)),
#                                 (Inches(5), Inches(3.5), Inches(3.5), Inches(0.6)),
#                                 (Inches(1), Inches(5.5), Inches(3.5), Inches(0.6)),
#                                 (Inches(5), Inches(5.5), Inches(3.5), Inches(0.6))
#                             ]
                    
#                     # Limit to maximum 4 images per slide
#                     images = images[:min(len(images), 4)]
                    
#                     # Add images to slide
#                     successfully_added_images = []
#                     for i, img_path in enumerate(images[:min(len(images), len(img_positions))]):
#                         if os.path.exists(img_path):
#                             left, top, width, _ = img_positions[i]
                            
#                             # Add the image
#                             slide.shapes.add_picture(img_path, left, top, width=width)
#                             successfully_added_images.append(img_path)
                    
#                     # Add a single consolidated caption for all images on the slide - ONLY if images were added
#                     if successfully_added_images:
#                         # Create a caption that combines information or uses a general description
#                         combined_caption = "Images: "
#                         image_captions = []
                        
#                         for img_path in successfully_added_images:
#                             if img_path in image_lookup and image_lookup[img_path].get("caption"):
#                                 image_captions.append(image_lookup[img_path]["caption"])
                        
#                         # If we have captions, create a consolidated one
#                         if image_captions:
#                             if len(image_captions) == 1:
#                                 combined_caption = image_captions[0]
#                             else:
#                                 # For multiple images, use a more intelligent consolidated caption
#                                 if len(image_captions) > 1:
#                                     try:
#                                         # Initialize the model
#                                         model = genai.GenerativeModel(GEMINI_MODEL)
                                        
#                                         # Create a prompt for generating a consolidated caption
#                                         prompt = f"""
#                                         Create a brief, cohesive caption (under 25 words) that describes this collection of images together.
#                                         Focus on their relationship, common themes, or progression.
                                        
#                                         Individual image captions:
#                                         {chr(10).join([f"- {caption}" for caption in image_captions])}
#                                         """
                                        
#                                         response = model.generate_content(
#                                             prompt,
#                                             generation_config={"temperature": 0.2, "max_output_tokens": 100}
#                                         )
                                        
#                                         combined_caption = response.text.strip()
                                        
#                                         # Add fallback in case of empty response
#                                         if not combined_caption:
#                                             combined_caption = "Collection of related images showing " + ", ".join(image_captions[:2]) + "..."
#                                     except Exception as e:
#                                         print(f"Error generating combined caption: {e}")
#                                         combined_caption = "Collection of related images for this slide"
                        
#                         # Add the combined caption below the images
#                         caption_box = slide.shapes.add_textbox(
#                             Inches(1), Inches(6.5), 
#                             Inches(8), Inches(1.0)  # Increased height to accommodate multiple lines
#                         )
#                         caption_frame = caption_box.text_frame
#                         caption_frame.text = combined_caption
                        
#                         # Enable word wrapping for multi-line captions
#                         caption_frame.word_wrap = True
                        
#                         # Style caption
#                         for paragraph in caption_frame.paragraphs:
#                             paragraph.alignment = PP_ALIGN.CENTER
#                             for run in paragraph.runs:
#                                 run.font.size = Pt(10)
#                                 run.font.italic = True
#                                 run.font.color.rgb = RGBColor(*text_rgb)
                
#                 # Add tables with captions if present in the slide (limit to 1 table per slide)
#                 tables = slide_content.get("tables", [])
#                 if tables:
#                     # Take only the first table
#                     table_path = tables[0]
                    
#                     if table_path in table_lookup:
#                         table_data = table_lookup[table_path]
                        
#                         # Get table DataFrame
#                         try:
#                             df = pd.read_csv(table_path)
                            
#                             # Create simplified table reference in slides
#                             rows, cols = df.shape
#                             max_rows = min(rows, 10)  # Limit to 10 rows max for presentation
#                             max_cols = min(cols, 6)   # Limit to 6 columns
                            
#                             # Position table
#                             left = Inches(1)
#                             top = Inches(2.5)
#                             width = Inches(8)
#                             height = Inches(3)
                            
#                             # Add table to slide
#                             shape_table = slide.shapes.add_table(
#                                 max_rows + 1,  # +1 for header
#                                 max_cols,
#                                 left, top, width, height
#                             )
                            
#                             tbl = shape_table.table
                            
#                             # Add headers
#                             for j, col_name in enumerate(df.columns[:max_cols]):
#                                 cell = tbl.cell(0, j)
#                                 cell.text = str(col_name)
#                                 # Style headers
#                                 for paragraph in cell.text_frame.paragraphs:
#                                     for run in paragraph.runs:
#                                         run.font.bold = True
#                                         run.font.color.rgb = RGBColor(*primary_rgb)
                            
#                             # Add data
#                             for row_idx in range(max_rows):
#                                 for col_idx in range(max_cols):
#                                     if row_idx < rows and col_idx < cols:
#                                         cell = tbl.cell(row_idx + 1, col_idx)  # +1 for header row
#                                         cell.text = str(df.iloc[row_idx, col_idx])
                            
#                             # Add table caption if available
#                             if table_data.get("caption"):
#                                 caption = table_data["caption"]
#                                 caption_box = slide.shapes.add_textbox(
#                                     Inches(1), Inches(6.5),  # Adjusted position similar to image captions
#                                     Inches(8), Inches(1.0)  # Increased height to accommodate multiple lines
#                                 )
#                                 caption_frame = caption_box.text_frame
#                                 caption_frame.text = f"Table: {caption}"
                                
#                                 # Enable word wrapping for multi-line captions
#                                 caption_frame.word_wrap = True
                                
#                                 # Style caption
#                                 for paragraph in caption_frame.paragraphs:
#                                     paragraph.alignment = PP_ALIGN.CENTER
#                                     for run in paragraph.runs:
#                                         run.font.size = Pt(10)
#                                         run.font.italic = True
#                                         run.font.color.rgb = RGBColor(*text_rgb)
                        
#                         except Exception as table_err:
#                             print(f"Error adding table to slide: {table_err}")
                
#                 # Add speaker notes if available
#                 if "notes" in slide_content and slide_content["notes"]:
#                     slide.notes_slide.notes_text_frame.text = slide_content["notes"]
            
#             # Save the presentation with timestamp
#             timestamp = time.strftime("%Y%m%d_%H%M%S")
#             output_filename = f"presentation_{timestamp}.pptx"
#             output_path = str(OUTPUT_DIR / output_filename)
#             prs.save(output_path)
            
#             return output_path
            
#         except Exception as e:
#             st.error(f"Error creating PowerPoint: {str(e)}")
#             return ""

# UI Components
st.title("PDF to PowerPoint Generator")
st.markdown("Upload a PDF and generate a PowerPoint presentation with AI assistance.")

# API Key input
with st.expander("API Key Configuration", expanded=not os.getenv("GEMINI_API_KEY")):
    api_key_input = st.text_input(
        "Gemini API Key",
        value=os.getenv("GEMINI_API_KEY", ""),
        type="password",
        key="api_key"
    )
    
    if st.button("Test API Connection"):
        if setup_gemini_api():
            st.success("API connection successful!")
        else:
            st.error("Failed to connect to API.")

# File upload
uploaded_file = st.file_uploader("Upload PDF", type=["pdf"])

# Presentation configuration
with st.expander("Presentation Settings", expanded=True):
    user_prompt = st.text_area(
        "Describe what you want in the presentation",
        "Create a concise, professional presentation that summarizes the key points in the document."
    )
    
    # Add presentation time input
    presentation_time = st.number_input(
        "Presentation Duration (minutes)",
        min_value=5,
        max_value=120,
        value=15,
        step=5,
        help="Specify how long the presentation will be. This helps optimize content amount and slide count."
    )
    
    # Color theme options
    st.subheader("Presentation Theme")
    col1, col2 = st.columns(2)
    
    with col1:
        primary_color = st.color_picker("Primary Color", "#3a86ff")
        text_color = st.color_picker("Text Color", "#333333")
    
    with col2:
        secondary_color = st.color_picker("Secondary Color", "#8338ec")
        background_color = st.color_picker("Background Color", "#ffffff")

# Process button
if uploaded_file and st.button("Generate Presentation"):
    if not setup_gemini_api():
        st.error("Please configure the Gemini API key first.")
    else:
        # Create progress bar
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        # Save uploaded file to temp location
        pdf_bytes = uploaded_file.read()
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            tmp_file.write(pdf_bytes)
            pdf_path = tmp_file.name
        
        try:
            # Step 1: Process PDF to extract content, images, and tables
            status_text.text("Step 1/3: Processing PDF...")
            document_data = process_pdf(pdf_path)
            progress_bar.progress(1/3)
            
            # Step 2: Optimize prompt based on presentation time and document structure
            status_text.text("Step 2/3: Optimizing presentation structure...")
            optimized_prompt = optimize_prompt_for_presentation_time(
                user_prompt, 
                document_data, 
                presentation_time
            )
            
            # Add user-selected theme information to prompt
            theme_prompt = f"""
            Use the following color theme for the presentation:
            - Primary color: {primary_color}
            - Secondary color: {secondary_color}
            - Text color: {text_color}
            - Background color: {background_color}
            """
            
            combined_prompt = optimized_prompt + "\n\n" + theme_prompt
            progress_bar.progress(2/3)
            
            # Step 3: Generate presentation content
            status_text.text("Step 3/3: Generating presentation content...")
            presentation_content = generate_presentation_content(
                document_data,
                combined_prompt
            )
            
            # Inject user-selected theme colors if not already present
            if "theme" not in presentation_content:
                presentation_content["theme"] = {}
            
            presentation_content["theme"]["primary_color"] = primary_color
            presentation_content["theme"]["secondary_color"] = secondary_color
            presentation_content["theme"]["text_color"] = text_color
            presentation_content["theme"]["background_color"] = background_color
            
            progress_bar.progress(3/3)
            
            # Create PowerPoint
            status_text.text("Creating PowerPoint presentation...")
            output_path = create_powerpoint(presentation_content, document_data)
            
            # Complete
            status_text.text("Presentation generation complete!")
            
            # Provide download link
            with open(output_path, "rb") as file:
                st.download_button(
                    label="Download Presentation",
                    data=file,
                    file_name=os.path.basename(output_path),
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            
            # Display some statistics
            num_images = len(document_data.get("images", []))
            st.success(f"Successfully generated presentation with {len(presentation_content.get('slides', []))} slides and {num_images} images.")
            
        except Exception as e:
            st.error(f"Error generating presentation: {str(e)}")
        
        finally:
            # Clean up
            try:
                os.unlink(pdf_path)
            except:
                pass

# Define a function to run PyTorch operations in a separate process
def run_pytorch_operation(func, *args, **kwargs):
    """Run a PyTorch operation in a separate process to avoid conflicts with Streamlit"""
    def wrapper(queue, func_name, *args, **kwargs):
        try:
            # Import PyTorch and other dependencies only in this process
            import docling
            # Run the function and put the result in the queue
            result = getattr(docling, func_name)(*args, **kwargs)
            queue.put(("success", result))
        except Exception as e:
            queue.put(("error", str(e)))
    
    # Create a queue for communication
    queue = multiprocessing.Queue()
    # Start a process
    process = multiprocessing.Process(
        target=wrapper, 
        args=(queue, func.__name__) + args, 
        kwargs=kwargs
    )
    process.start()
    process.join()
    
    # Get the result
    status, result = queue.get()
    if status == "error":
        raise RuntimeError(f"Error in PyTorch operation: {result}")
    return result

# Then use this function for heavy PyTorch operations:
# document = run_pytorch_operation(docling.Document.from_pdf, pdf_path)

# Create cached functions for expensive PyTorch operations
@st.cache_data
def process_pdf_cached(pdf_path):
    """Process PDF with PyTorch-based operations, with results cached"""
    doc = docling.Document.from_pdf(pdf_path)
    # Extract all the data you need from the document
    result = {
        "pages": [],
        "text": ""
    }
    
    # Process each page
    for page in doc.pages:
        page_data = {
            "text": page.extract_text(),
            "tables": [table.to_dict() for table in page.extract_tables()],
            "images": [img.to_dict() for img in page.extract_images()]
        }
        result["pages"].append(page_data)
        result["text"] += page_data["text"] + "\n\n"
    
    return result
